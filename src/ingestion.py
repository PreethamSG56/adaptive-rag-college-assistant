"""
Document Loader and Chunker
============================
Supports: PDF, DOCX, PPTX, TXT, PNG, JPG files.
Chunks text with configurable size and overlap.

Speed optimisations:
  - EasyOCR reader is a module-level singleton (loads once per process)
  - OpenCV image preprocessing before OCR (grayscale → denoise → threshold)
  - PDF pages are OCR'd in parallel via ThreadPoolExecutor
  - NLP post-processing skipped for high-quality digital text

NLP Pipeline (for handwritten / scanned notes):
  OCR (EasyOCR)  →  Noise removal  →  Spell-correction (SymSpell)
  →  Sentence segmentation (spaCy)  →  Keyword extraction
  →  Confidence scoring  →  Chunk enrichment
"""

from __future__ import annotations

import hashlib
import os
import pickle
import re
import warnings
from concurrent.futures import ThreadPoolExecutor, as_completed
from dataclasses import dataclass, field
from pathlib import Path
from typing import List, Optional

# ── EasyOCR singleton — load the model ONCE for the entire process ──────────
_easyocr_reader = None

def _get_ocr_reader():
    """Lazy-load EasyOCR reader once and cache it globally."""
    global _easyocr_reader
    if _easyocr_reader is not None:
        return _easyocr_reader
    try:
        import easyocr
        print("[OCR] Loading EasyOCR model (one-time, subsequent calls are instant)...")
        with warnings.catch_warnings():
            warnings.simplefilter("ignore")
            _easyocr_reader = easyocr.Reader(
                ['en'],
                gpu=False,
                verbose=False,
                model_storage_directory=None,  # use default cache
            )
        print("[OCR] EasyOCR model loaded and cached.")
    except Exception as e:
        print(f"[OCR] Failed to load EasyOCR: {e}")
        _easyocr_reader = None
    return _easyocr_reader


# ── OpenCV image preprocessing — makes OCR faster AND more accurate ─────────
def _preprocess_image(image_bytes: bytes) -> bytes:
    """
    Preprocess image before OCR:
      1. Decode bytes → numpy array
      2. Convert to grayscale (removes colour noise)
      3. Resize to 2x if small (EasyOCR accuracy improves at higher resolution)
      4. Apply fast non-local means denoise
      5. Adaptive threshold → clean black-on-white binary image
      6. Re-encode as PNG bytes

    Returns preprocessed PNG bytes, or original bytes if OpenCV unavailable.
    """
    try:
        import cv2
        import numpy as np

        # Decode
        arr = np.frombuffer(image_bytes, dtype=np.uint8)
        img = cv2.imdecode(arr, cv2.IMREAD_COLOR)
        if img is None:
            return image_bytes

        # Grayscale
        gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)

        # Upscale if the image is small (improves OCR speed by reducing failed reads)
        h, w = gray.shape
        if max(h, w) < 1000:
            scale = 2.0
            gray = cv2.resize(gray, (int(w * scale), int(h * scale)),
                              interpolation=cv2.INTER_CUBIC)

        # Fast denoise (h=10 for handwritten, 5 for printed)
        gray = cv2.fastNlMeansDenoising(gray, h=10, templateWindowSize=7, searchWindowSize=21)

        # Adaptive threshold — turns handwriting into clean black-on-white
        binary = cv2.adaptiveThreshold(
            gray, 255,
            cv2.ADAPTIVE_THRESH_GAUSSIAN_C,
            cv2.THRESH_BINARY,
            blockSize=31,
            C=10,
        )

        # Re-encode
        _, buf = cv2.imencode('.png', binary)
        return buf.tobytes()
    except Exception as e:
        print(f"[OCR] OpenCV preprocessing skipped: {e}")
        return image_bytes


def _transcribe_image_local(image_bytes: bytes, run_nlp: bool = True) -> tuple:
    """
    Fast OCR pipeline:
      1. OpenCV preprocessing (grayscale, denoise, threshold)
      2. EasyOCR via cached singleton reader (no model reload!)
      3. NLP post-processing (only for low-confidence / handwritten pages)

    Returns
    -------
    tuple (cleaned_text: str, ocr_confidence: float, keywords: list[str])
    """
    # Step 1 — preprocess image for faster, cleaner OCR
    processed_bytes = _preprocess_image(image_bytes)

    # Step 2 — OCR with cached reader (loads only once per process)
    reader = _get_ocr_reader()
    if reader is None:
        return "", 0.0, []

    raw_text = ""
    try:
        print("[OCR] Scanning image...")
        results = reader.readtext(
            processed_bytes,
            detail=0,
            paragraph=True,     # group lines → fewer calls, faster
            batch_size=8,       # process 8 text regions at once
        )
        raw_text = " ".join(results).strip()
        print(f"[OCR] Extracted {len(raw_text)} chars")
    except Exception as e:
        print(f"[OCR] EasyOCR error: {e}")
        return "", 0.0, []

    if not raw_text:
        return "", 0.0, []

    # Step 3 — NLP post-processing
    if run_nlp:
        try:
            from src.nlp_processor import process_ocr_text, _ocr_quality_score
            # Quick confidence check — skip heavy NLP for clean text
            quick_conf = _ocr_quality_score(raw_text)
            run_spell = quick_conf < 0.85   # only spell-correct if text looks noisy
            nlp_result = process_ocr_text(
                raw_text,
                run_spell_correction=run_spell,
                run_spacy=True,
                is_handwritten=True,
            )
            print(
                f"[NLP] conf={nlp_result.ocr_confidence:.2f} "
                f"spell={nlp_result.spell_corrected} "
                f"kw={len(nlp_result.keywords)}"
            )
            return nlp_result.cleaned_text, nlp_result.ocr_confidence, nlp_result.keywords
        except Exception as e:
            print(f"[NLP] Post-processing error: {e} — using raw OCR text")

    return raw_text, 0.5, []


@dataclass
class Document:
    text: str
    source: str
    page: int = 0
    chunk_id: int = 0
    metadata: dict = field(default_factory=dict)


def _load_txt(path: Path) -> str:
    """Load plain text or markdown file."""
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        return f.read()


def _ocr_page(args: tuple) -> tuple:
    """Worker function for parallel PDF page OCR. Returns (page_index, text, conf, keywords)."""
    i, img_bytes = args
    ocr_text, conf, kws = _transcribe_image_local(img_bytes)
    return i, ocr_text, conf, kws


def _load_pdf(path: Path) -> tuple:
    """
    Extract text from PDF using PyMuPDF.
    - Digital text pages: extracted instantly (no OCR).
    - Scanned/handwritten pages: OCR'd in PARALLEL threads for maximum speed.

    Returns
    -------
    tuple (text: str, page_meta: list[dict])
        page_meta contains per-page ocr_confidence and keywords.
    """
    page_meta: list[dict] = []

    # Try PyMuPDF (fitz) first
    try:
        import fitz
        doc = fitz.open(path)
        total_pages = len(doc)

        # ── Pass 1: extract digital text & identify scanned pages ──────────
        page_texts   = [""] * total_pages
        page_confs   = [1.0] * total_pages
        page_kws     = [[] for _ in range(total_pages)]
        ocr_jobs     = []   # (page_index, img_bytes) for scanned pages

        for i, page in enumerate(doc):
            text = page.get_text("text").strip()
            if len(text) >= 150:
                # Digital text — fast, no OCR needed
                page_texts[i] = text
            else:
                # Scanned/handwritten — render at 2x zoom and queue for OCR
                print(f"[Loader] Page {i+1}/{total_pages} → queued for OCR")
                pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                ocr_jobs.append((i, pix.tobytes("png")))

        # ── Pass 2: parallel OCR for scanned pages ─────────────────────────
        if ocr_jobs:
            max_workers = min(4, len(ocr_jobs))   # cap at 4 threads (CPU-bound)
            print(f"[Loader] Running OCR on {len(ocr_jobs)} page(s) with {max_workers} thread(s)...")
            with ThreadPoolExecutor(max_workers=max_workers) as pool:
                futures = {pool.submit(_ocr_page, job): job[0] for job in ocr_jobs}
                for future in as_completed(futures):
                    pg_i, ocr_text, conf, kws = future.result()
                    if ocr_text:
                        page_texts[pg_i] = ocr_text
                        page_confs[pg_i] = conf
                        page_kws[pg_i]   = kws
                        print(f"[Loader] Page {pg_i+1} OCR done — conf={conf:.2f} {len(ocr_text)}chars")

        # ── Assemble results ───────────────────────────────────────────────
        pages = []
        for i in range(total_pages):
            if page_texts[i]:
                pages.append(page_texts[i])
                page_meta.append({
                    "page": i + 1,
                    "ocr_confidence": page_confs[i],
                    "keywords": page_kws[i],
                })

        full_text = "\n\n".join(pages)
        print(f"[Loader] PDF done: {len(full_text)} chars, {len(pages)} pages "
              f"({len(ocr_jobs)} OCR'd in parallel)")
        return full_text, page_meta
    except Exception as e:
        print(f"[Loader] PyMuPDF error {path.name}: {e}")

    # Fallback to pdfplumber (no OCR meta)
    try:
        import pdfplumber
        with pdfplumber.open(path) as pdf:
            pages = [page.extract_text() or "" for page in pdf.pages]
        full_text = "\n\n".join(pages)
        print(f"[Loader] PDF loaded via pdfplumber fallback: {len(full_text)} chars")
        return full_text, []
    except Exception as e:
        print(f"[Loader] pdfplumber error {path.name}: {e}")

    # Final fallback to pypdf
    try:
        from pypdf import PdfReader
        reader = PdfReader(str(path))
        pages = [page.extract_text() or "" for page in reader.pages]
        full_text = "\n\n".join(pages)
        print(f"[Loader] PDF loaded via pypdf fallback: {len(full_text)} chars")
        return full_text, []
    except Exception as e:
        print(f"[Loader] pypdf error {path.name}: {e}")
        return "", []


def _load_docx(path: Path) -> str:
    """Extract text from DOCX file."""
    try:
        from docx import Document as DocxDoc
        doc = DocxDoc(str(path))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        return "\n\n".join(paragraphs)
    except Exception as e:
        print(f"[Loader] DOCX error {path.name}: {e}")
        return ""


def _load_pptx(path: Path) -> str:
    """Extract text from PPTX file."""
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        slides = []
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            if texts:
                slides.append(f"[Slide {i+1}]\n" + "\n".join(texts))
        return "\n\n".join(slides)
    except Exception as e:
        print(f"[Loader] PPTX error {path.name}: {e}")
        return ""


def _load_image(path: Path) -> tuple:
    """
    Extract handwritten or printed text from an image using offline EasyOCR + NLP.

    Returns
    -------
    tuple (text: str, ocr_confidence: float, keywords: list[str])
    """
    try:
        with open(path, "rb") as f:
            img_bytes = f.read()
        print(f"[Loader] Running OCR + NLP on image {path.name}...")
        return _transcribe_image_local(img_bytes)
    except Exception as e:
        print(f"[Loader] Image OCR error {path.name}: {e}")
        return "", 0.0, []


# ── Fast regex keyword extractor (no spaCy load needed) ─────────────────────
_CAPITAL_WORD = re.compile(r"\b[A-Z][a-z]{2,}(?:\s+[A-Z][a-z]{2,}){0,2}\b")
_TECH_TERM    = re.compile(r"\b[A-Za-z][a-z]*(?:[A-Z][a-z]+)+\b")   # camelCase

def _fast_keywords(text: str, max_kw: int = 20) -> list:
    """
    Extract keywords using regex only — no spaCy import needed.
    ~100x faster than spaCy for large files.
    """
    caps   = _CAPITAL_WORD.findall(text[:8000])
    camel  = _TECH_TERM.findall(text[:8000])
    seen, result = set(), []
    for kw in caps + camel:
        kl = kw.lower()
        if kl not in seen and len(kw) > 3:
            seen.add(kl)
            result.append(kw)
        if len(result) >= max_kw:
            break
    return result


def load_file(path: Path) -> tuple:
    """
    Dispatch to the correct loader based on file extension.

    Returns
    -------
    tuple (text: str, ocr_confidence: float, keywords: list[str], page_meta: list[dict])
    """
    ext = path.suffix.lower()
    if ext in (".txt", ".md"):
        raw = _load_txt(path)
        # Fast regex keywords — avoids slow spaCy startup
        return raw, 1.0, _fast_keywords(raw), []
    elif ext == ".pdf":
        text, page_meta = _load_pdf(path)
        all_kws: list = []
        for pm in page_meta:
            all_kws.extend(pm.get("keywords", []))
        avg_conf = (
            sum(pm.get("ocr_confidence", 1.0) for pm in page_meta) / len(page_meta)
            if page_meta else 1.0
        )
        return text, avg_conf, list(dict.fromkeys(all_kws)), page_meta
    elif ext == ".docx":
        raw = _load_docx(path)
        return raw, 1.0, _fast_keywords(raw), []
    elif ext in (".pptx", ".ppt"):
        raw = _load_pptx(path)
        return raw, 1.0, _fast_keywords(raw), []
    elif ext in (".png", ".jpg", ".jpeg"):
        text, conf, kws = _load_image(path)
        return text, conf, kws, []
    else:
        return "", 1.0, [], []


def chunk_text(
    text: str,
    chunk_size: int = 500,
    overlap: int = 50,
) -> List[str]:
    """
    Split text into overlapping character-level chunks.

    The function first collapses redundant whitespace so that chunk sizes
    are predictable, then slices the text with the requested ``overlap``
    to preserve context at boundaries.

    Parameters
    ----------
    text : str
        Cleaned input text.
    chunk_size : int
        Maximum number of characters per chunk (default 500).
    overlap : int
        Number of characters to repeat at the start of each successive
        chunk so that retrieval is not confused by arbitrary boundaries
        (default 50).

    Returns
    -------
    List[str]
        List of non-empty text chunks.
    """
    # Collapse all whitespace runs to a single space
    text = re.sub(r"[\t ]+", " ", text)          # horizontal whitespace
    text = re.sub(r"\n{3,}", "\n\n", text)        # max two consecutive newlines
    text = text.strip()
    if not text:
        return []

    chunks: list[str] = []
    start = 0
    while start < len(text):
        end = min(start + chunk_size, len(text))
        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)
        if end == len(text):
            break
        start += chunk_size - overlap
    return chunks


# ── File hash cache (skip unchanged files on re-index) ──────────────────────
import hashlib

def _file_hash(path: Path) -> str:
    """Fast MD5 of file content — used to detect unchanged files."""
    h = hashlib.md5()
    with open(path, "rb") as f:
        for chunk in iter(lambda: f.read(65536), b""):
            h.update(chunk)
    return h.hexdigest()


def _process_single_file(
    file_path: Path,
    chunk_size: int,
    overlap: int,
) -> List[Document]:
    """Load + chunk one file. Used as the parallel worker."""
    print(f"[Loader] ▶ {file_path.name}")
    text, ocr_confidence, keywords, _ = load_file(file_path)
    if not text.strip():
        print(f"[Loader] ✗ Empty: {file_path.name}")
        return []

    # Append keyword footer once (fast, no spaCy)
    kw_footer = f"\n[KEY_CONCEPTS: {', '.join(keywords[:15])}]" if keywords else ""

    chunks = chunk_text(text, chunk_size, overlap)
    docs = []
    for i, chunk in enumerate(chunks):
        docs.append(Document(
            text=chunk + (kw_footer if i == 0 else ""),  # footer only on first chunk
            source=file_path.name,
            chunk_id=i,
            metadata={
                "file": str(file_path),
                "ocr_confidence": ocr_confidence,
                "keywords": ",".join(keywords[:10]),
                "is_handwritten": ocr_confidence < 0.95,
            },
        ))
    print(f"[Loader] ✓ {file_path.name} → {len(chunks)} chunks (conf={ocr_confidence:.2f})")
    return docs


def ingest_documents(
    source: str,
    chunk_size: int = 500,
    overlap: int = 50,
    max_workers: int = 4,
    cache_file: Optional[str] = None,
) -> List[Document]:
    """
    Load all supported files from a file path or directory.

    Speed features:
    - Parallel file processing (ThreadPoolExecutor, default 4 workers)
    - Hash-based cache: unchanged files are skipped entirely
    - Fast regex keyword extraction (no spaCy startup cost)

    Returns a flat list of Document objects.
    """
    source_path = Path(source)

    if source_path.is_dir():
        files = (
            list(source_path.rglob("*.txt"))
            + list(source_path.rglob("*.md"))
            + list(source_path.rglob("*.pdf"))
            + list(source_path.rglob("*.docx"))
            + list(source_path.rglob("*.pptx"))
            + list(source_path.rglob("*.png"))
            + list(source_path.rglob("*.jpg"))
            + list(source_path.rglob("*.jpeg"))
        )
    elif source_path.is_file():
        files = [source_path]
    else:
        print(f"[Loader] Path not found: {source}")
        return []

    if not files:
        return []

    # ── Hash-based incremental cache ─────────────────────────────────────────
    cache_path = Path(cache_file) if cache_file else source_path.parent / ".ingest_cache.pkl"
    hash_cache: dict = {}
    if cache_path.exists():
        try:
            with open(cache_path, "rb") as f:
                hash_cache = pickle.load(f)
        except Exception:
            hash_cache = {}

    import pickle as _pkl

    new_hashes: dict = {}
    files_to_process: List[Path] = []
    cached_docs: List[Document] = []

    for fp in files:
        fh = _file_hash(fp)
        new_hashes[str(fp)] = fh
        if hash_cache.get(str(fp)) == fh and str(fp) in hash_cache.get("_docs", {}):
            # File unchanged — reuse cached documents
            cached_docs.extend(hash_cache["_docs"][str(fp)])
            print(f"[Loader] ⚡ Cached (unchanged): {fp.name}")
        else:
            files_to_process.append(fp)

    # ── Parallel processing of new/changed files ──────────────────────────────
    fresh_docs: List[Document] = []
    if files_to_process:
        workers = min(max_workers, len(files_to_process))
        print(f"[Loader] Processing {len(files_to_process)} file(s) with {workers} worker(s)...")
        with ThreadPoolExecutor(max_workers=workers) as pool:
            futures = {
                pool.submit(_process_single_file, fp, chunk_size, overlap): fp
                for fp in files_to_process
            }
            for future in as_completed(futures):
                try:
                    result = future.result()
                    fresh_docs.extend(result)
                    # Cache this file's docs
                    fp = futures[future]
                    if "_docs" not in new_hashes:
                        new_hashes["_docs"] = {}
                    new_hashes["_docs"][str(fp)] = result
                except Exception as e:
                    print(f"[Loader] Error processing {futures[future].name}: {e}")

    # ── Save updated hash cache ───────────────────────────────────────────────
    if files_to_process:
        # Merge old cached docs into new_hashes
        if "_docs" not in new_hashes:
            new_hashes["_docs"] = {}
        for fp_str, docs in hash_cache.get("_docs", {}).items():
            if fp_str not in new_hashes["_docs"]:
                new_hashes["_docs"][fp_str] = docs
        try:
            with open(cache_path, "wb") as f:
                _pkl.dump(new_hashes, f)
        except Exception as e:
            print(f"[Loader] Cache save failed: {e}")

    all_docs = cached_docs + fresh_docs
    print(f"[Loader] Total: {len(all_docs)} chunks "
          f"({len(cached_docs)} from cache, {len(fresh_docs)} fresh)")
    return all_docs
